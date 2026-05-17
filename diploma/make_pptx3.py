# -*- coding: utf-8 -*-
"""
Diploma defense — modern dark Bento Grid presentation
10 x 7.5 in (4:3), 19 slides
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ──────────────────────────────────────────────────────────────────────────────
# CANVAS
# ──────────────────────────────────────────────────────────────────────────────
W = Inches(10)
H = Inches(7.5)
MX = Inches(0.45)       # horizontal margin
CW = W - MX * 2        # 9.10 in usable width
TY = Inches(0.18)       # title-area top
DDY = Inches(0.9)       # divider below title
CY = Inches(1.05)       # content area start
CE = Inches(7.0)        # content area end
CH = CE - CY            # 5.95 in content height
GAP = Inches(0.18)      # gap between cards
CW2 = (CW - GAP) // 2  # 2-col card width ≈ 4.46 in
CW3 = (CW - GAP * 2) // 3  # 3-col card width ≈ 2.91 in

# ──────────────────────────────────────────────────────────────────────────────
# PALETTE
# ──────────────────────────────────────────────────────────────────────────────
BG   = RGBColor(0x09, 0x13, 0x24)   # near-black navy
CARD = RGBColor(0x10, 0x1E, 0x38)   # card surface
CARD2= RGBColor(0x17, 0x2A, 0x4C)   # lighter card variant
A1   = RGBColor(0x3D, 0x8E, 0xFF)   # electric blue (primary)
A2   = RGBColor(0x06, 0xD6, 0xD6)   # cyan
A3   = RGBColor(0xFB, 0xBF, 0x24)   # amber
A4   = RGBColor(0x34, 0xD3, 0x99)   # green
A5   = RGBColor(0xA7, 0x8B, 0xFA)   # violet
T1   = RGBColor(0xF1, 0xF5, 0xF9)   # primary text
T2   = RGBColor(0x94, 0xA3, 0xB8)   # secondary text
T3   = RGBColor(0x47, 0x55, 0x6B)   # muted text
LN   = RGBColor(0x1E, 0x35, 0x58)   # divider line

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BL = prs.slide_layouts[6]


# ──────────────────────────────────────────────────────────────────────────────
# CORE PRIMITIVES
# ──────────────────────────────────────────────────────────────────────────────
def R(sl, x, y, w, h, fill=None, lc=None):
    s = sl.shapes.add_shape(1, int(x), int(y), int(w), int(h))
    s.line.fill.background()
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if lc:
        s.line.color.rgb = lc
        s.line.width = Pt(0.75)
    else:
        s.line.fill.background()
    return s


def T(sl, text, x, y, w, h, size=13, bold=False, color=None,
      align=PP_ALIGN.LEFT, italic=False):
    bx = sl.shapes.add_textbox(int(x), int(y), int(w), int(h))
    tf = bx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color or T1
    return bx


def TM(sl, rows, x, y, w, h, align=PP_ALIGN.LEFT):
    """Multi-paragraph text.  rows = [(text, size, bold, color), ...]"""
    bx = sl.shapes.add_textbox(int(x), int(y), int(w), int(h))
    tf = bx.text_frame
    tf.word_wrap = True
    for i, (text, size, bold, color) in enumerate(rows):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return bx


# ──────────────────────────────────────────────────────────────────────────────
# SLIDE SCAFFOLDING
# ──────────────────────────────────────────────────────────────────────────────
def bg_dark(sl):
    R(sl, 0, 0, W, H, fill=BG)


def top_stripe(sl, color=A1, thickness=Inches(0.05)):
    R(sl, 0, 0, W, thickness, fill=color)


def bot_stripe(sl, color=A1):
    R(sl, 0, H - Inches(0.05), W, Inches(0.05), fill=color)


def pg_divider(sl, y=DDY):
    R(sl, MX, y, CW, Inches(0.022), fill=LN)


def pg_num(sl, n, total=19):
    T(sl, f"{n} / {total}", W - Inches(1.4), H - Inches(0.38),
      Inches(1.2), Inches(0.28), size=10, color=T3, align=PP_ALIGN.RIGHT)


def page_header(sl, title, tag=None, tag_color=A1):
    if tag:
        T(sl, tag.upper(), MX, TY, CW, Inches(0.28),
          size=10, bold=True, color=tag_color)
        T(sl, title, MX, Inches(0.44), CW, Inches(0.56),
          size=32, bold=True, color=T1)
    else:
        T(sl, title, MX, TY, CW, Inches(0.7),
          size=32, bold=True, color=T1)
    pg_divider(sl)


def std(n, title, tag=None, accent=A1):
    sl = prs.slides.add_slide(BL)
    bg_dark(sl)
    top_stripe(sl, accent)
    bot_stripe(sl, accent)
    page_header(sl, title, tag, accent)
    pg_num(sl, n)
    return sl


# ──────────────────────────────────────────────────────────────────────────────
# CARD HELPERS
# ──────────────────────────────────────────────────────────────────────────────
def accent_card(sl, x, y, w, h, title, body, accent=A1,
                ts=15, bs=12, fill=CARD):
    """Card with left accent stripe, title, and body text."""
    R(sl, x, y, w, h, fill=fill)
    R(sl, x, y, Inches(0.07), h, fill=accent)
    px = x + Inches(0.18)
    pw = w - Inches(0.25)
    if title:
        T(sl, title, px, y + Inches(0.14), pw, Inches(0.36),
          size=ts, bold=True, color=T1)
    if body:
        by = y + Inches(0.54) if title else y + Inches(0.14)
        T(sl, body, px, by, pw, h - (by - y) - Inches(0.12),
          size=bs, color=T2)


def num_card(sl, x, y, w, h, num, title, body=None, accent=A1, ts=14, bs=12):
    """Card with circle number badge."""
    R(sl, x, y, w, h, fill=CARD)
    badge = Inches(0.44)
    by2 = y + (h - badge) // 2
    R(sl, x + Inches(0.16), by2, badge, badge, fill=accent)
    T(sl, str(num), x + Inches(0.16), by2, badge, badge,
      size=14, bold=True, color=T1, align=PP_ALIGN.CENTER)
    tx = x + Inches(0.16) + badge + Inches(0.14)
    tw = w - (tx - x) - Inches(0.1)
    T(sl, title, tx, y + Inches(0.1), tw, Inches(0.32),
      size=ts, bold=True, color=T1)
    if body:
        T(sl, body, tx, y + Inches(0.44), tw, h - Inches(0.5),
          size=bs, color=T2)


def icon_stat(sl, x, y, w, h, value, label, accent=A1):
    """Large statistic card."""
    R(sl, x, y, w, h, fill=CARD)
    R(sl, x, y, w, Inches(0.06), fill=accent)
    T(sl, value, x, y + Inches(0.1), w, Inches(0.72),
      size=40, bold=True, color=accent, align=PP_ALIGN.CENTER)
    T(sl, label, x + Inches(0.1), y + Inches(0.85), w - Inches(0.2), Inches(0.4),
      size=12, color=T2, align=PP_ALIGN.CENTER)


def pill(sl, x, y, w, h, text, fill=CARD2, text_color=T1, ts=11):
    R(sl, x, y, w, h, fill=fill)
    T(sl, text, x, y, w, h, size=ts, bold=True, color=text_color,
      align=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────────────────────────────────────
# PHONE MOCKUP
# ──────────────────────────────────────────────────────────────────────────────
PW = Inches(3.2)   # phone width
PH = Inches(5.65)  # phone height


def phone(sl, px, py):
    """Draw phone frame. Returns (sx, sy, sw, sh) — inner app screen area."""
    # Body
    R(sl, px, py, PW, PH, fill=RGBColor(0x0C, 0x11, 0x1E))
    # Notch
    nw = Inches(0.66)
    R(sl, px + (PW - nw) // 2, py + Inches(0.05), nw, Inches(0.17),
      fill=RGBColor(0x05, 0x08, 0x12))
    # Volume / power buttons
    R(sl, px - Inches(0.04), py + Inches(0.82), Inches(0.04), Inches(0.3),
      fill=RGBColor(0x16, 0x22, 0x3A))
    R(sl, px + PW, py + Inches(0.72), Inches(0.04), Inches(0.45),
      fill=RGBColor(0x16, 0x22, 0x3A))
    # Screen
    sx = px + Inches(0.08)
    sy = py + Inches(0.26)
    sw = PW - Inches(0.16)
    sh = PH - Inches(0.5)
    R(sl, sx, sy, sw, sh, fill=RGBColor(0x08, 0x14, 0x2C))
    # Status bar
    R(sl, sx, sy, sw, Inches(0.22), fill=RGBColor(0x04, 0x0A, 0x18))
    # Bottom nav bar
    nav_y = sy + sh - Inches(0.4)
    R(sl, sx, nav_y, sw, Inches(0.4), fill=RGBColor(0x04, 0x0A, 0x18))
    for i in range(3):
        dot_x = sx + Inches(0.32 + i * 0.84)
        dot_c = A1 if i == 0 else T3
        R(sl, dot_x, nav_y + Inches(0.13), Inches(0.16), Inches(0.14), fill=dot_c)
    # Home bar
    R(sl, px + (PW - Inches(0.48)) // 2, py + PH - Inches(0.14),
      Inches(0.48), Inches(0.06), fill=T3)
    return sx, sy + Inches(0.22), sw, sh - Inches(0.62)


def app_bar(sl, sx, sy, sw, text, fill=RGBColor(0x07, 0x14, 0x2A)):
    R(sl, sx, sy, sw, Inches(0.3), fill=fill)
    T(sl, text, sx + Inches(0.08), sy + Inches(0.04), sw - Inches(0.14), Inches(0.22),
      size=9, bold=True, color=T1)


def app_row(sl, sx, sy, sw, title, sub="", accent=A1, h=Inches(0.42)):
    R(sl, sx + Inches(0.04), sy, sw - Inches(0.08), h,
      fill=RGBColor(0x0D, 0x1C, 0x38))
    R(sl, sx + Inches(0.04), sy, Inches(0.05), h, fill=accent)
    T(sl, title, sx + Inches(0.13), sy + Inches(0.04),
      sw - Inches(0.22), Inches(0.2), size=8, bold=True, color=T1)
    if sub:
        T(sl, sub, sx + Inches(0.13), sy + Inches(0.26),
          sw - Inches(0.22), Inches(0.15), size=7, color=T3)


def app_field(sl, sx, sy, sw, label, h=Inches(0.32)):
    R(sl, sx + Inches(0.04), sy, sw - Inches(0.08), h,
      fill=RGBColor(0x0D, 0x1C, 0x38), lc=RGBColor(0x1A, 0x38, 0x6A))
    T(sl, label, sx + Inches(0.1), sy + Inches(0.06),
      sw - Inches(0.2), Inches(0.2), size=8, color=T3)


def app_btn(sl, sx, sy, sw, text, accent=A1):
    R(sl, sx + Inches(0.04), sy, sw - Inches(0.08), Inches(0.32), fill=accent)
    T(sl, text, sx + Inches(0.04), sy + Inches(0.04), sw - Inches(0.08), Inches(0.24),
      size=9, bold=True, color=T1, align=PP_ALIGN.CENTER)


def app_badge(sl, sx, sy, text, fill=A4):
    bw = Inches(0.58)
    R(sl, sx, sy, bw, Inches(0.18), fill=fill)
    T(sl, text, sx, sy, bw, Inches(0.18), size=7, bold=True,
      color=T1, align=PP_ALIGN.CENTER)


# ──────────────────────────────────────────────────────────────────────────────
# SCREEN SLIDE HELPER
# ──────────────────────────────────────────────────────────────────────────────
def screen_slide(n, caption, subtitle, features, draw_fn, accent=A1):
    sl = std(n, "Экранные формы", "ЭКРАННЫЕ ФОРМЫ", accent)
    # Screen title inside content
    T(sl, caption, MX, CY, Inches(4.7), Inches(0.38),
      size=18, bold=True, color=accent)
    T(sl, subtitle, MX, CY + Inches(0.4), Inches(4.7), Inches(0.28),
      size=11, color=T3)
    # Feature cards
    fh = (CH - Inches(0.75) - GAP * (len(features) - 1)) // len(features)
    for i, (ft, fb) in enumerate(features):
        fy = CY + Inches(0.75) + i * (fh + GAP)
        R(sl, MX, fy, Inches(4.5), fh, fill=CARD)
        R(sl, MX, fy, Inches(0.06), fh, fill=accent)
        # dot badge
        bd = Inches(0.3)
        by2 = fy + (fh - bd) // 2
        R(sl, MX + Inches(0.12), by2, bd, bd, fill=CARD2)
        T(sl, str(i + 1), MX + Inches(0.12), by2, bd, bd,
          size=10, bold=True, color=accent, align=PP_ALIGN.CENTER)
        T(sl, ft, MX + Inches(0.54), fy + Inches(0.08), Inches(3.85), Inches(0.28),
          size=12, bold=True, color=T1)
        T(sl, fb, MX + Inches(0.54), fy + Inches(0.36), Inches(3.85),
          fh - Inches(0.43), size=11, color=T2)
    # Phone
    phone_x = Inches(5.2)
    phone_y = CY - Inches(0.06)
    psx, psy, psw, psh = phone(sl, phone_x, phone_y)
    draw_fn(sl, psx, psy, psw, psh)
    # Caption below phone
    T(sl, caption, phone_x, phone_y + PH + Inches(0.08), PW, Inches(0.28),
      size=11, color=T3, align=PP_ALIGN.CENTER)
    return sl


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BL)
bg_dark(sl)
top_stripe(sl, A1)
bot_stripe(sl, A1)

# Gradient-like header band
R(sl, 0, Inches(0.05), W, Inches(1.1), fill=RGBColor(0x0D, 0x1F, 0x42))

# University line
T(sl, "Северо-Кавказская государственная академия  ·  Институт цифровых технологий  ·  09.03.03 — Прикладная информатика",
  MX, Inches(0.14), CW, Inches(0.3), size=10, color=T3, align=PP_ALIGN.CENTER)

# ВКР tag
pill(sl, Inches(3.8), Inches(0.55), Inches(2.4), Inches(0.3),
     "ВЫПУСКНАЯ КВАЛИФИКАЦИОННАЯ РАБОТА", fill=A1,
     text_color=T1, ts=10)

# Big title
T(sl, "Разработка мобильного приложения для рассылки новостей и образовательных ресурсов высшего учебного заведения",
  MX, Inches(1.35), CW, Inches(1.35), size=26, bold=True, color=T1,
  align=PP_ALIGN.CENTER)

# Subtitle
T(sl, "с использованием технологии React Native",
  MX, Inches(2.68), CW, Inches(0.36), size=16, color=A1, align=PP_ALIGN.CENTER)

# Divider
R(sl, Inches(2.5), Inches(3.2), Inches(5.0), Inches(0.022), fill=LN)

# Tech pills row
pills = [("React Native", A1), ("Node.js", A4), ("PostgreSQL", A3),
         ("JWT Auth", A5), ("Expo SDK 50", A2)]
pw2 = Inches(1.72)
gap2 = Inches(0.12)
start_x = MX + (CW - len(pills) * pw2 - (len(pills) - 1) * gap2) // 2
for i, (txt, clr) in enumerate(pills):
    R(sl, start_x + i * (pw2 + gap2), Inches(3.38), pw2, Inches(0.3), fill=CARD2)
    R(sl, start_x + i * (pw2 + gap2), Inches(3.38), Inches(0.05), Inches(0.3), fill=clr)
    T(sl, txt, start_x + i * (pw2 + gap2) + Inches(0.1), Inches(3.38),
      pw2 - Inches(0.1), Inches(0.3), size=11, bold=True, color=clr)

# Student block
R(sl, MX, Inches(3.88), Inches(5.5), Inches(2.48), fill=CARD)
R(sl, MX, Inches(3.88), Inches(0.07), Inches(2.48), fill=A1)
TM(sl, [
    ("ВЫПОЛНИЛ:", 10, True, A1),
    ("Иванов Муса Насирович", 20, True, T1),
    ("Обучающийся ОФО, направление подготовки 09.03.03", 12, False, T2),
], MX + Inches(0.18), Inches(3.98), Inches(5.1), Inches(0.85))
TM(sl, [
    ("РУКОВОДИТЕЛЬ:", 10, True, A1),
    ("уч. степень, уч. звание, ФИО (полностью)", 13, False, T2),
], MX + Inches(0.18), Inches(4.95), Inches(5.1), Inches(0.65))
T(sl, "Черкесск, 2026", MX + Inches(0.18), Inches(5.72), Inches(5.1), Inches(0.4),
  size=14, bold=True, color=A3)

# Right decorative block
R(sl, Inches(6.1), Inches(3.88), Inches(3.45), Inches(2.48), fill=CARD2)
R(sl, Inches(6.1), Inches(3.88), Inches(3.45), Inches(0.05), fill=A2)
T(sl, "СТЕК ТЕХНОЛОГИЙ", Inches(6.2), Inches(3.96), Inches(3.2), Inches(0.28),
  size=11, bold=True, color=A2)
stack_items = [
    ("▸ React Native 0.73 + Expo SDK 50", T1),
    ("▸ Node.js 20 + Express + TypeScript 5", T1),
    ("▸ PostgreSQL 16 + Prisma ORM 5", T1),
    ("▸ JWT + bcrypt + Zustand + Axios", T2),
    ("▸ Expo SecureStore (Keychain/Keystore)", T2),
]
for i, (txt, clr) in enumerate(stack_items):
    T(sl, txt, Inches(6.2), Inches(4.3 + i * 0.4), Inches(3.2), Inches(0.38),
      size=12, color=clr)

pg_num(sl, 1)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — АКТУАЛЬНОСТЬ
# ══════════════════════════════════════════════════════════════════════════════
sl = std(2, "Актуальность", "АКТУАЛЬНОСТЬ", A1)

# 2x2 bento grid
rh = (CH - GAP * 2) // 3         # 3 rows; top two are bento, bottom is conclusion
rh2 = (CH - GAP) // 2            # for 2-row use

problems = [
    (A1,  "Разрозненность источников",
     "Сайт, email, мессенджеры — информация рассеяна по независимым каналам без единой точки входа"),
    (A3,  "Отсутствие проактивности",
     "Студент сам ищет новости. Нет гарантии охвата аудитории; нет уведомлений об обновлениях"),
    (A5,  "Нет разграничения ролей",
     "Существующие решения не разделяют права ADMIN / TEACHER / STUDENT на публикацию контента"),
    (A2,  "Избыточность аналогов",
     "MS Teams, Moodle — высокие лицензионные издержки и функциональность, избыточная для задачи"),
]
h_row = (CH - GAP * 3 - Inches(1.0)) // 2
for i, (clr, ttl, bdy) in enumerate(problems):
    col = i % 2
    row = i // 2
    x = MX + col * (CW2 + GAP)
    y = CY + row * (h_row + GAP)
    accent_card(sl, x, y, CW2, h_row, ttl, bdy, accent=clr, ts=14, bs=12)

# Full-width conclusion card
cy2 = CY + 2 * (h_row + GAP)
R(sl, MX, cy2, CW, Inches(0.9), fill=RGBColor(0x0D, 0x25, 0x48))
R(sl, MX, cy2, CW, Inches(0.06), fill=A1)
T(sl, "ВЫВОД:", MX + Inches(0.18), cy2 + Inches(0.1), Inches(1.0), Inches(0.28),
  size=11, bold=True, color=A1)
T(sl, "Требуется специализированное мобильное приложение — единый интерфейс публикации институциональных новостей и образовательных ресурсов с ролевым разграничением доступа.",
  MX + Inches(1.15), cy2 + Inches(0.1), CW - Inches(1.25), Inches(0.72),
  size=13, color=T1)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — ОБЪЕКТ / ПРЕДМЕТ / ЦЕЛЬ
# ══════════════════════════════════════════════════════════════════════════════
sl = std(3, "Объект, предмет и цель", "ПОСТАНОВКА ЗАДАЧИ", A2)

items3 = [
    (A2,  "Объект исследования",
     "Информационная среда высшего учебного заведения как совокупность информационных потоков между участниками образовательного процесса.", Inches(1.55)),
    (A5,  "Предмет исследования",
     "Методы и средства разработки мобильного приложения для централизованной публикации и распространения институциональных новостей и образовательных ресурсов.", Inches(1.55)),
    (A1,  "Цель",
     "Разработать кроссплатформенное мобильное приложение для централизованной публикации новостей и управления образовательными ресурсами вуза с разграничением прав по ролям STUDENT, TEACHER и ADMIN.", Inches(2.2)),
]
y3 = CY
for clr, ttl, bdy, h3 in items3:
    accent_card(sl, MX, y3, CW, h3, ttl, bdy, accent=clr, ts=16, bs=13, fill=CARD)
    y3 += h3 + GAP


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — ЗАДАЧИ
# ══════════════════════════════════════════════════════════════════════════════
sl = std(4, "Задачи исследования", "ЗАДАЧИ", A3)

tasks = [
    ("Анализ предметной области и информационных потоков вуза", ""),
    ("Сравнительный анализ программных аналогов", "MS Teams, Moodle Mobile, ВКонтакте"),
    ("Проектирование информационной модели", "4 сущности, схема БД в 3НФ, REST API контракт"),
    ("Обоснование выбора технологического стека", "React Native, Node.js, PostgreSQL, Prisma"),
    ("Реализация серверной части", "Node.js + Express + TypeScript + Prisma ORM + JWT"),
    ("Реализация мобильного клиента", "React Native + Expo SDK 50 + Zustand + SecureStore"),
    ("Контрольный пример", "Три сценария: ADMIN, TEACHER, STUDENT"),
]
# 4 left + 3 right, slightly different heights
card_h = (CH - GAP * 3) // 4
for i in range(4):
    num_card(sl, MX, CY + i * (card_h + GAP), CW2, card_h,
             i + 1, tasks[i][0], tasks[i][1] if tasks[i][1] else None,
             accent=A3, ts=13, bs=11)
card_h2 = (CH - GAP * 2) // 3
for i in range(3):
    num_card(sl, MX + CW2 + GAP, CY + i * (card_h2 + GAP), CW2, card_h2,
             i + 5, tasks[i + 4][0], tasks[i + 4][1] if tasks[i + 4][1] else None,
             accent=A3, ts=13, bs=11)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — СТЕК ТЕХНОЛОГИЙ
# ══════════════════════════════════════════════════════════════════════════════
sl = std(5, "Стек технологий", "ТЕХНОЛОГИИ", A4)

layers = [
    ("Мобильный клиент", A1, [
        ("React Native 0.73 + Expo SDK 50",
         "Единая кодовая база iOS/Android, нативные API, Expo managed workflow"),
        ("React Navigation 6 (Stack + BottomTabs)",
         "Многоуровневая навигация с независимой историей вкладок"),
        ("Zustand",
         "Минималистичный global store: user, role, флаги аутентификации"),
        ("Axios + interceptors",
         "Автоматический Bearer-токен, прозрачный refresh при HTTP 401"),
        ("Expo SecureStore",
         "Keychain (iOS) / Keystore (Android) — шифрование токенов на уровне ОС"),
    ]),
    ("Серверная часть + БД", A2, [
        ("Node.js 20 + Express 4 + TypeScript 5",
         "Типобезопасная модульная серверная логика, ESM modules"),
        ("Prisma ORM 5",
         "Type-safe запросы, автоматические миграции, генерация TS-типов из схемы"),
        ("PostgreSQL 16",
         "Реляционная СУБД, enum-типы, 3НФ, FK-ограничения целостности"),
        ("jsonwebtoken",
         "Access token (24 ч) + Refresh token (7 сут), подпись серверным секретом"),
        ("bcrypt  (cost factor 10)",
         "Адаптивное хеширование паролей, хранение паролей в открытом виде исключено"),
    ]),
]
lw = (CW - GAP) // 2
for li, (ltitle, lcolor, items) in enumerate(layers):
    lx = MX + li * (lw + GAP)
    R(sl, lx, CY, lw, Inches(0.4), fill=lcolor)
    T(sl, ltitle, lx + Inches(0.12), CY + Inches(0.06), lw - Inches(0.2),
      Inches(0.28), size=14, bold=True, color=T1)
    ih = (CH - Inches(0.4) - GAP * 4) // 5
    for ii, (tech, desc) in enumerate(items):
        iy = CY + Inches(0.4) + GAP + ii * (ih + GAP)
        R(sl, lx, iy, lw, ih, fill=CARD)
        R(sl, lx, iy, Inches(0.06), ih, fill=lcolor)
        T(sl, tech, lx + Inches(0.14), iy + Inches(0.06), lw - Inches(0.22), Inches(0.3),
          size=12, bold=True, color=T1)
        T(sl, desc, lx + Inches(0.14), iy + Inches(0.38), lw - Inches(0.22),
          ih - Inches(0.44), size=11, color=T2)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — БИЗНЕС-ПРОЦЕССЫ
# ══════════════════════════════════════════════════════════════════════════════
sl = std(6, "Бизнес-процессы системы", "ФУНКЦИОНАЛЬНОСТЬ", A5)

procs = [
    (A1,  "1. Аутентификация пользователя",
     "Форма → POST /api/auth/login → JWT-пара → SecureStore → role в Zustand → навигация"),
    (A3,  "2. Публикация новости (ADMIN)",
     "AdminNewsFormScreen → POST /api/news + Bearer → authenticate + authorize → INSERT"),
    (A2,  "3. Просмотр ленты новостей",
     "GET /api/news → FlatList → фильтр по categoryId → переход на NewsDetailScreen"),
    (A4,  "4. Создание ресурса (TEACHER / ADMIN)",
     "Форма → POST /api/resources + Bearer → authorize(['TEACHER','ADMIN']) → INSERT"),
    (A5,  "5. Каталог ресурсов со студентом",
     "GET /api/resources → цветовая кодировка типов LINK/FILE/TEXT → детальный просмотр"),
    (A1,  "6. Управление категориями (ADMIN)",
     "CRUD /api/categories → категории используются в фильтрах обоих разделов"),
    (A3,  "7. Обновление токенов",
     "Axios interceptor: 401 → POST /api/auth/refresh → новая пара → повтор запроса"),
]
ph = (CH - GAP * 6) // 7
for i, (clr, ttl, bdy) in enumerate(procs):
    accent_card(sl, MX, CY + i * (ph + GAP), CW, ph, ttl, bdy,
                accent=clr, ts=13, bs=11)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — СХЕМА ДАННЫХ
# ══════════════════════════════════════════════════════════════════════════════
sl = std(7, "Схема данных", "ИНФОРМАЦИОННАЯ МОДЕЛЬ", A2)

EW = Inches(4.28)
EH = Inches(2.62)
entities = [
    (MX,            CY,           A1,  "users",
     ["id  INT  PK  AI", "fullName  VARCHAR(255)  NOT NULL",
      "email  VARCHAR(255)  UNIQUE  NOT NULL", "password  VARCHAR(255)  [bcrypt hash]",
      "role  ENUM(ADMIN, TEACHER, STUDENT)"]),
    (MX + EW + GAP, CY,           A2,  "categories",
     ["id  INT  PK  AI", "name  VARCHAR(100)  NOT NULL",
      "type  ENUM(NEWS, RESOURCE, COMMON)", "createdAt / updatedAt  TIMESTAMP"]),
    (MX,            CY + EH + GAP, A3, "news",
     ["id · title · shortDescription · content",
      "imageUrl  NULL  |  published  BOOL  DEFAULT false",
      "publishedAt  TIMESTAMP  NULL",
      "authorId  → users.id  ON DELETE RESTRICT",
      "categoryId  → categories.id  ON DELETE SET NULL"]),
    (MX + EW + GAP, CY + EH + GAP, A4, "resources",
     ["id · title · description · type  ENUM(LINK,FILE,TEXT)",
      "fileUrl  NULL  |  textContent  TEXT  NULL",
      "authorId  → users.id  ON DELETE RESTRICT",
      "categoryId  → categories.id  ON DELETE SET NULL"]),
]
for ex, ey, clr, name, fields in entities:
    R(sl, ex, ey, EW, EH, fill=CARD)
    R(sl, ex, ey, EW, Inches(0.42), fill=clr)
    T(sl, name, ex + Inches(0.12), ey + Inches(0.06), EW - Inches(0.2), Inches(0.3),
      size=15, bold=True, color=T1)
    for fi, f in enumerate(fields):
        T(sl, f, ex + Inches(0.12), ey + Inches(0.5 + fi * 0.4),
          EW - Inches(0.2), Inches(0.38), size=10.5, color=T2)

# Relationship lines
# users → news (vertical)
mid_x = MX + EW * 0.35
R(sl, mid_x, CY + EH, Inches(0.022), GAP, fill=T3)
T(sl, "1:N", mid_x + Inches(0.05), CY + EH + GAP * 0.1,
  Inches(0.6), Inches(0.28), size=10, bold=True, color=A1)
# users → resources (horizontal then down)
R(sl, MX + EW * 0.65, CY + EH, Inches(0.022), GAP, fill=T3)
T(sl, "1:N", MX + EW * 0.65 + Inches(0.05), CY + EH + GAP * 0.1,
  Inches(0.6), Inches(0.28), size=10, bold=True, color=A1)
# categories → news
mid_x2 = MX + EW + GAP + EW * 0.3
R(sl, mid_x2, CY + EH, Inches(0.022), GAP, fill=T3)
T(sl, "0..N", mid_x2 + Inches(0.05), CY + EH + GAP * 0.1,
  Inches(0.7), Inches(0.28), size=10, bold=True, color=A2)
# categories → resources
mid_x3 = MX + EW + GAP + EW * 0.65
R(sl, mid_x3, CY + EH, Inches(0.022), GAP, fill=T3)
T(sl, "0..N", mid_x3 + Inches(0.05), CY + EH + GAP * 0.1,
  Inches(0.7), Inches(0.28), size=10, bold=True, color=A2)

# Note at bottom
T(sl, "ON DELETE RESTRICT (authorId)  ·  ON DELETE SET NULL (categoryId)  ·  Схема в третьей нормальной форме",
  MX, Inches(7.1), CW, Inches(0.26), size=10, color=T3, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDES 8–16 — ЭКРАННЫЕ ФОРМЫ
# ══════════════════════════════════════════════════════════════════════════════

# ── Slide 8: Login ───────────────────────────────────────────────────────────
def draw_login(sl, sx, sy, sw, sh):
    app_bar(sl, sx, sy, sw, "Вход в систему")
    sy += Inches(0.3)
    R(sl, sx, sy, sw, Inches(0.5), fill=RGBColor(0x05, 0x12, 0x28))
    T(sl, "University App", sx, sy + Inches(0.1), sw, Inches(0.3),
      size=13, bold=True, color=A1, align=PP_ALIGN.CENTER)
    sy += Inches(0.58)
    app_field(sl, sx, sy, sw, "Email адрес")
    app_field(sl, sx, sy + Inches(0.38), sw, "Пароль")
    app_btn(sl, sx, sy + Inches(0.84), sw, "Войти", A1)
    T(sl, "Регистрация", sx, sy + Inches(1.24), sw, Inches(0.22),
      size=8, color=A1, align=PP_ALIGN.CENTER)

screen_slide(8, "1. Экран входа в систему",
             "Единственный экран для неаутентифицированного пользователя",
             [("POST /api/auth/login", "email + пароль → bcrypt.compare → JWT-пара (access 24ч, refresh 7сут)"),
              ("Expo SecureStore", "Токены шифруются Keychain (iOS) / Keystore (Android)"),
              ("Zustand store", "user, role → формирует набор вкладок навигатора"),
              ("Автоматическая инициализация", "При запуске — чтение токенов из SecureStore → восстановление сессии")],
             draw_login, A1)

# ── Slide 9: Home ────────────────────────────────────────────────────────────
def draw_home(sl, sx, sy, sw, sh):
    app_bar(sl, sx, sy, sw, "Главная")
    sy += Inches(0.34)
    T(sl, "Последние новости", sx + Inches(0.06), sy, sw, Inches(0.22),
      size=8, bold=True, color=T3)
    sy += Inches(0.26)
    for t, sub in [("Изменение расписания сессии", "14 мая 2026 · ADMIN"),
                   ("Итоги весенней олимпиады", "10 мая 2026 · ADMIN")]:
        app_row(sl, sx, sy, sw, t, sub, A1, Inches(0.44))
        sy += Inches(0.5)
    T(sl, "Ресурсы", sx + Inches(0.06), sy, sw, Inches(0.22),
      size=8, bold=True, color=T3)
    sy += Inches(0.26)
    for t, tp, clr in [("Лекция №3: Основы БД", "LINK", A1),
                       ("Методичка по Python", "FILE", A3)]:
        app_row(sl, sx, sy, sw, t, tp, clr, Inches(0.4))
        sy += Inches(0.46)

screen_slide(9, "2. Главный экран",
             "Агрегированная лента: 3 новости + 3 ресурса",
             [("GET /api/news + GET /api/resources", "Параллельные запросы (limit=3), без дополнительных запросов"),
              ("BottomTabNavigator", "«Главная», «Новости», «Ресурсы» + «Управление» только для ADMIN"),
              ("Адаптивный UI по роли", "Набор вкладок и кнопок зависит от role в Zustand store"),
              ("Переходы к полным спискам", "Кнопки «Все новости» / «Все ресурсы» на экране")],
             draw_home, A1)

# ── Slide 10: News List ──────────────────────────────────────────────────────
def draw_newslist(sl, sx, sy, sw, sh):
    app_bar(sl, sx, sy, sw, "Новости")
    sy += Inches(0.34)
    for i, cat in enumerate(["Все", "Адм.", "Учеб."]):
        fc = A1 if i == 0 else RGBColor(0x0D, 0x1C, 0x38)
        R(sl, sx + Inches(0.04 + i * 0.62), sy, Inches(0.56), Inches(0.22), fill=fc)
        T(sl, cat, sx + Inches(0.04 + i * 0.62), sy, Inches(0.56), Inches(0.22),
          size=7, bold=True, color=T1, align=PP_ALIGN.CENTER)
    sy += Inches(0.28)
    for t, d in [("Изменение расписания", "14 мая"),
                 ("Итоги олимпиады", "10 мая"),
                 ("Набор в кружки", "5 мая"),
                 ("Конференция 2026", "1 мая")]:
        app_row(sl, sx, sy, sw, t, d, A1, Inches(0.4))
        sy += Inches(0.46)

screen_slide(10, "3. Список новостей",
             "Лента опубликованных новостей с фильтрацией",
             [("GET /api/news + GET /api/categories?type=NEWS", "Параллельная загрузка новостей и категорий при монтировании"),
              ("FlatList с pull-to-refresh", "Обновление при потягивании вниз, lazy-loading карточек"),
              ("Фильтры по категориям", "GET /api/news?categoryId=N при выборе — горизонтальный скролл"),
              ("Переход к детальному просмотру", "Stack-навигатор → NewsDetailScreen при нажатии на карточку")],
             draw_newslist, A1)

# ── Slide 11: News Detail ────────────────────────────────────────────────────
def draw_newsdetail(sl, sx, sy, sw, sh):
    app_bar(sl, sx, sy, sw, "Детали новости")
    sy += Inches(0.3)
    R(sl, sx, sy, sw, Inches(0.62), fill=RGBColor(0x06, 0x12, 0x2A))
    T(sl, "[изображение]", sx, sy + Inches(0.2), sw, Inches(0.22),
      size=8, color=T3, align=PP_ALIGN.CENTER)
    sy += Inches(0.66)
    R(sl, sx + Inches(0.04), sy, Inches(0.9), Inches(0.18),
      fill=RGBColor(0x0D, 0x25, 0x4C))
    T(sl, "Адм. объявления", sx + Inches(0.06), sy + Inches(0.02),
      Inches(0.9), Inches(0.14), size=6, color=A1)
    sy += Inches(0.24)
    T(sl, "Изменение расписания на период сессии",
      sx + Inches(0.06), sy, sw - Inches(0.1), Inches(0.34),
      size=8, bold=True, color=T1)
    sy += Inches(0.38)
    T(sl, "Администратор  ·  14 мая 2026",
      sx + Inches(0.06), sy, sw - Inches(0.1), Inches(0.18), size=7, color=T3)
    sy += Inches(0.24)
    R(sl, sx + Inches(0.04), sy, sw - Inches(0.08), Inches(0.36),
      fill=RGBColor(0x0A, 0x20, 0x44))
    R(sl, sx + Inches(0.04), sy, Inches(0.04), Inches(0.36), fill=A1)
    T(sl, "В связи с началом сессии...", sx + Inches(0.1), sy + Inches(0.06),
      sw - Inches(0.2), Inches(0.26), size=7, color=RGBColor(0xBB, 0xD4, 0xFF))

screen_slide(11, "4. Детальный просмотр новости",
             "Полная страница публикации с изображением и текстом",
             [("GET /api/news/:id", "Полный объект с вложенными author и category объектами"),
              ("Категориальный бейдж + автор + дата", "Метаданные публикации в едином блоке заголовка"),
              ("Выделенный блок краткого описания", "Левая синяя полоса визуально отделяет вводную часть"),
              ("Полный текст без ограничений длины", "ScrollView — материал любого объёма читается удобно")],
             draw_newsdetail, A1)

# ── Slide 12: Resources List ─────────────────────────────────────────────────
def draw_reslist(sl, sx, sy, sw, sh):
    app_bar(sl, sx, sy, sw, "Ресурсы")
    sy += Inches(0.34)
    for i, cat in enumerate(["Все", "Учеб.", "Доп."]):
        fc = A2 if i == 0 else RGBColor(0x0D, 0x1C, 0x38)
        R(sl, sx + Inches(0.04 + i * 0.62), sy, Inches(0.56), Inches(0.22), fill=fc)
        T(sl, cat, sx + Inches(0.04 + i * 0.62), sy, Inches(0.56), Inches(0.22),
          size=7, bold=True, color=T1, align=PP_ALIGN.CENTER)
    sy += Inches(0.28)
    for t, tp, clr in [("Лекция №3: Основы реляционных БД", "LINK", A1),
                       ("Методичка по Python", "FILE", A3),
                       ("Краткий конспект лекции", "TEXT", A4)]:
        R(sl, sx + Inches(0.04), sy, sw - Inches(0.08), Inches(0.44),
          fill=RGBColor(0x0D, 0x1C, 0x38))
        R(sl, sx + Inches(0.04), sy, Inches(0.05), Inches(0.44), fill=clr)
        R(sl, sx + sw - Inches(0.6), sy + Inches(0.04), Inches(0.52), Inches(0.18), fill=clr)
        T(sl, tp, sx + sw - Inches(0.6), sy + Inches(0.04), Inches(0.52), Inches(0.18),
          size=7, bold=True, color=T1, align=PP_ALIGN.CENTER)
        T(sl, t, sx + Inches(0.12), sy + Inches(0.06), sw - Inches(0.78), Inches(0.18),
          size=7, bold=True, color=T1)
        T(sl, "Иванов И.И.  ·  12 мая", sx + Inches(0.12), sy + Inches(0.26),
          sw - Inches(0.22), Inches(0.15), size=7, color=T3)
        sy += Inches(0.5)

screen_slide(12, "5. Каталог образовательных ресурсов",
             "Список материалов с типами и фильтрацией по категориям",
             [("Цветовая кодировка типов", "LINK — синий  ·  FILE — оранжевый  ·  TEXT — зелёный"),
              ("Кнопка «Добавить» — TEACHER / ADMIN", "Кнопка в шапке скрыта от пользователей с ролью STUDENT"),
              ("GET /api/resources?categoryId=N", "Фильтрация при выборе категории, горизонтальный скролл"),
              ("Карточка: тип, название, автор, дата", "Вся необходимая информация без перехода в детали")],
             draw_reslist, A2)

# ── Slide 13: Resource Detail ────────────────────────────────────────────────
def draw_resdetail(sl, sx, sy, sw, sh):
    app_bar(sl, sx, sy, sw, "Ресурс")
    sy += Inches(0.3)
    R(sl, sx + Inches(0.04), sy, sw - Inches(0.08), Inches(0.22),
      fill=RGBColor(0x08, 0x1E, 0x3E))
    T(sl, "LINK  ·  Учебные материалы", sx + Inches(0.1), sy + Inches(0.03),
      sw - Inches(0.2), Inches(0.16), size=7, bold=True, color=A1)
    sy += Inches(0.28)
    T(sl, "Лекция №3: Основы реляционных баз данных",
      sx + Inches(0.06), sy, sw - Inches(0.1), Inches(0.38),
      size=8, bold=True, color=T1)
    sy += Inches(0.42)
    T(sl, "Конспект лекции по теме нормализации отношений, первая, вторая и третья нормальные формы...",
      sx + Inches(0.06), sy, sw - Inches(0.1), Inches(0.4), size=7, color=T2)
    sy += Inches(0.48)
    T(sl, "Иванов И.И.  ·  12 мая 2026",
      sx + Inches(0.06), sy, sw - Inches(0.1), Inches(0.18), size=7, color=T3)
    sy += Inches(0.28)
    app_btn(sl, sx, sy, sw, "Открыть ссылку", A1)
    sy += Inches(0.38)
    app_btn(sl, sx, sy, sw, "Поделиться", CARD2)

screen_slide(13, "6. Детальный просмотр ресурса",
             "Содержимое адаптируется к типу LINK / FILE / TEXT",
             [("GET /api/resources/:id", "Объект с полями fileUrl или textContent в зависимости от типа"),
              ("LINK: «Открыть ссылку»", "WebBrowser.openBrowserAsync() — встроенный браузер"),
              ("FILE: кнопка загрузки", "Ссылка на файл, доступная для скачивания"),
              ("TEXT: текстовое содержимое", "textContent хранится непосредственно в БД, рендерится инлайн")],
             draw_resdetail, A2)

# ── Slide 14: Admin News List ────────────────────────────────────────────────
def draw_adminnews(sl, sx, sy, sw, sh):
    app_bar(sl, sx, sy, sw, "Новости  [ADMIN]")
    sy += Inches(0.32)
    R(sl, sx + sw - Inches(0.62), sy - Inches(0.28), Inches(0.58), Inches(0.24), fill=A1)
    T(sl, "+ Создать", sx + sw - Inches(0.62), sy - Inches(0.28),
      Inches(0.58), Inches(0.24), size=7, bold=True, color=T1, align=PP_ALIGN.CENTER)
    for t, pub in [("Изменение расписания на период сессии", True),
                   ("Итоги весенней олимпиады", True),
                   ("Черновик: анонс конференции", False)]:
        R(sl, sx + Inches(0.04), sy, sw - Inches(0.08), Inches(0.5),
          fill=RGBColor(0x0D, 0x1C, 0x38))
        badge_c = A4 if pub else A3
        badge_t = "Опубл." if pub else "Черновик"
        app_badge(sl, sx + sw - Inches(0.66), sy + Inches(0.04), badge_t, badge_c)
        T(sl, t, sx + Inches(0.1), sy + Inches(0.05), sw - Inches(0.76),
          Inches(0.22), size=7, bold=True, color=T1)
        R(sl, sx + Inches(0.1), sy + Inches(0.3), Inches(0.28), Inches(0.14), fill=A1)
        T(sl, "Ред.", sx + Inches(0.1), sy + Inches(0.3), Inches(0.28), Inches(0.14),
          size=6, bold=True, color=T1, align=PP_ALIGN.CENTER)
        R(sl, sx + Inches(0.44), sy + Inches(0.3), Inches(0.28), Inches(0.14), fill=A5)
        T(sl, "Удал.", sx + Inches(0.44), sy + Inches(0.3), Inches(0.28), Inches(0.14),
          size=6, bold=True, color=T1, align=PP_ALIGN.CENTER)
        sy += Inches(0.56)

screen_slide(14, "7. Административный список новостей",
             "Полный список включая черновики — только для ADMIN",
             [("authorize(['ADMIN']) на сервере", "Middleware проверяет роль из JWT — клиентская защита лишь UX"),
              ("Индикаторы статуса Published/Draft", "Администратор сразу видит, что опубликовано, что черновик"),
              ("Кнопки Редактировать + Удалить", "PUT /api/news/:id  ·  DELETE /api/news/:id"),
              ("Кнопка «+ Создать» в шапке", "Переход на AdminNewsFormScreen с пустой формой")],
             draw_adminnews, A5)

# ── Slide 15: Create News Form ───────────────────────────────────────────────
def draw_newsform(sl, sx, sy, sw, sh):
    app_bar(sl, sx, sy, sw, "Создать новость")
    sy += Inches(0.32)
    for lbl in ["Заголовок *", "Краткое описание *", "Полный текст *", "URL изображения"]:
        app_field(sl, sx, sy, sw, lbl)
        sy += Inches(0.36)
    R(sl, sx + Inches(0.04), sy, sw - Inches(0.08), Inches(0.28),
      fill=RGBColor(0x0D, 0x1C, 0x38), lc=RGBColor(0x1A, 0x38, 0x6A))
    T(sl, "Категория ▾", sx + Inches(0.1), sy + Inches(0.06),
      sw - Inches(0.2), Inches(0.18), size=8, color=T3)
    sy += Inches(0.36)
    app_btn(sl, sx, sy, sw, "Сохранить", A1)

screen_slide(15, "8. Форма создания и редактирования новости",
             "Единая форма для POST (создание) и PUT (редактирование)",
             [("Обязательные поля с валидацией", "Заголовок, краткое описание, полный текст — проверка перед запросом"),
              ("Выпадающий список категорий", "GET /api/categories?type=NEWS — категории загружаются динамически"),
              ("Режим редактирования", "GET /api/news/:id → предзаполнение формы текущими значениями"),
              ("POST или PUT в зависимости от режима", "Единый компонент — два сценария использования")],
             draw_newsform, A5)

# ── Slide 16: Categories ─────────────────────────────────────────────────────
def draw_categories(sl, sx, sy, sw, sh):
    app_bar(sl, sx, sy, sw, "Категории  [ADMIN]")
    sy += Inches(0.32)
    R(sl, sx + sw - Inches(0.62), sy - Inches(0.28), Inches(0.58), Inches(0.24), fill=A1)
    T(sl, "+ Добавить", sx + sw - Inches(0.62), sy - Inches(0.28),
      Inches(0.58), Inches(0.24), size=7, bold=True, color=T1, align=PP_ALIGN.CENTER)
    T(sl, "NEWS", sx + Inches(0.08), sy, sw, Inches(0.2),
      size=8, bold=True, color=A3)
    sy += Inches(0.24)
    for cat in ["Административные объявления", "Учебный процесс"]:
        R(sl, sx + Inches(0.04), sy, sw - Inches(0.08), Inches(0.3),
          fill=RGBColor(0x0D, 0x1C, 0x38))
        T(sl, cat, sx + Inches(0.1), sy + Inches(0.06), sw - Inches(0.5), Inches(0.18),
          size=7, color=T1)
        R(sl, sx + sw - Inches(0.5), sy + Inches(0.06), Inches(0.18), Inches(0.18), fill=A1)
        R(sl, sx + sw - Inches(0.26), sy + Inches(0.06), Inches(0.18), Inches(0.18), fill=A5)
        sy += Inches(0.36)
    T(sl, "RESOURCE", sx + Inches(0.08), sy + Inches(0.06), sw, Inches(0.2),
      size=8, bold=True, color=A2)
    sy += Inches(0.3)
    for cat in ["Учебные материалы", "Дополнительно"]:
        R(sl, sx + Inches(0.04), sy, sw - Inches(0.08), Inches(0.3),
          fill=RGBColor(0x0D, 0x1C, 0x38))
        T(sl, cat, sx + Inches(0.1), sy + Inches(0.06), sw - Inches(0.5), Inches(0.18),
          size=7, color=T1)
        R(sl, sx + sw - Inches(0.5), sy + Inches(0.06), Inches(0.18), Inches(0.18), fill=A1)
        R(sl, sx + sw - Inches(0.26), sy + Inches(0.06), Inches(0.18), Inches(0.18), fill=A5)
        sy += Inches(0.36)

screen_slide(16, "9. Экран управления категориями",
             "Классификаторы для обоих разделов — только ADMIN",
             [("Два раздела: NEWS и RESOURCE", "Категория NEWS не появится в фильтрах раздела «Ресурсы»"),
              ("CRUD /api/categories", "POST создать · PUT переименовать · DELETE удалить"),
              ("ON DELETE SET NULL", "Удаление категории не удаляет публикации — categoryId → NULL"),
              ("Влияет на фильтры обоих разделов", "Любое изменение сразу отражается в горизонтальных фильтрах")],
             draw_categories, A2)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — ВЫВОДЫ
# ══════════════════════════════════════════════════════════════════════════════
sl = std(17, "Выводы", "РЕЗУЛЬТАТЫ", A4)

conclusions = [
    (A1, "Анализ предметной области",
     "Выявлены ограничения разрозненных каналов коммуникации, обоснована необходимость специализированного решения"),
    (A3, "Сравнительный анализ аналогов",
     "MS Teams, Moodle, ВКонтакте: ни одно из решений не обеспечивает оптимального сочетания функций при контролируемых затратах"),
    (A2, "Информационная модель + REST API",
     "4 сущности в 3НФ (User, News, Resource, Category), 17 REST-эндпоинтов, ограничения ссылочной целостности"),
    (A5, "Серверная часть",
     "Node.js/Express/TypeScript/Prisma с JWT-аутентификацией, двухуровневым RBAC и bcrypt-хешированием паролей"),
    (A4, "Мобильный клиент",
     "11 экранов на React Native, адаптивный UI для трёх ролей, interceptor обновления токенов, Expo SecureStore"),
    (A1, "Контрольный пример",
     "Три сценария подтверждены. 403 при попытке обойти RBAC. Время ответа API — менее 200 мс"),
]
cw6 = (CW - GAP * 2) // 3
ch6 = (CH - GAP) // 2
for i, (clr, ttl, bdy) in enumerate(conclusions):
    col = i % 3
    row = i // 3
    x = MX + col * (cw6 + GAP)
    y = CY + row * (ch6 + GAP)
    accent_card(sl, x, y, cw6, ch6, ttl, bdy, accent=clr, ts=13, bs=11)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — ДОКЛАДЧИК
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BL)
bg_dark(sl)
top_stripe(sl, A1)
bot_stripe(sl, A1)
pg_num(sl, 18)

# Photo placeholder
R(sl, MX, Inches(0.6), Inches(2.9), Inches(3.6), fill=CARD)
R(sl, MX, Inches(0.6), Inches(2.9), Inches(0.06), fill=A1)
T(sl, "ФОТО", MX, Inches(2.1), Inches(2.9), Inches(0.5),
  size=16, color=T3, align=PP_ALIGN.CENTER)

# Name + info
T(sl, "Иванов Муса Насирович",
  Inches(3.65), Inches(0.65), Inches(5.9), Inches(0.58),
  size=24, bold=True, color=T1)
R(sl, Inches(3.65), Inches(1.28), Inches(5.9), Inches(0.022), fill=A1)
TM(sl, [
    ("Обучающийся очной формы обучения", 13, False, T2),
    ("09.03.03 — Прикладная информатика", 13, False, T2),
    ("Северо-Кавказская государственная академия", 13, True, T1),
], Inches(3.65), Inches(1.36), Inches(5.9), Inches(0.72))

T(sl, "+7 (___) ___-__-__", Inches(3.65), Inches(2.2), Inches(5.9), Inches(0.4),
  size=16, bold=True, color=T1)
T(sl, "email@ncsa.ru", Inches(3.65), Inches(2.62), Inches(5.9), Inches(0.4),
  size=16, bold=True, color=A1)

# University block
R(sl, MX, Inches(4.5), CW, Inches(2.65), fill=CARD)
R(sl, MX, Inches(4.5), CW, Inches(0.06), fill=A1)
T(sl, "СКГА", MX + Inches(0.2), Inches(4.62), Inches(3.5), Inches(1.15),
  size=72, bold=True, color=A1)
T(sl, "Северо-Кавказская\nгосударственная академия",
  Inches(3.8), Inches(4.75), Inches(5.8), Inches(0.8),
  size=18, bold=True, color=T1)
T(sl, "ncsa.ru", Inches(3.8), Inches(5.6), Inches(5.8), Inches(0.35),
  size=14, color=A1)
T(sl, "Черкесск, 2026", MX + Inches(0.2), Inches(5.7), Inches(3.5), Inches(0.38),
  size=14, color=T3)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — ФИНАЛЬНЫЙ
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BL)
bg_dark(sl)
top_stripe(sl, A1)
bot_stripe(sl, A1)

R(sl, 0, Inches(0.05), W, Inches(0.8), fill=RGBColor(0x0D, 0x1E, 0x40))
T(sl, "Северо-Кавказская государственная академия  ·  2026",
  MX, Inches(0.15), CW, Inches(0.35), size=11, color=T3, align=PP_ALIGN.CENTER)

T(sl, "Доклад окончен.",
  MX, Inches(1.3), CW, Inches(0.82), size=48, bold=True, color=A1,
  align=PP_ALIGN.CENTER)
T(sl, "Спасибо за внимание!",
  MX, Inches(2.15), CW, Inches(0.82), size=48, bold=True, color=T1,
  align=PP_ALIGN.CENTER)

R(sl, Inches(1.5), Inches(3.1), Inches(7.0), Inches(0.022), fill=LN)

# 4 stat cards
stats = [
    ("17",    "эндпоинтов REST API",  A1),
    ("11",    "экранов приложения",   A2),
    ("4",     "сущности в базе данных", A3),
    ("<200мс","время ответа API",     A4),
]
sw2 = (CW - GAP * 3) // 4
for i, (val, lbl, clr) in enumerate(stats):
    sx2 = MX + i * (sw2 + GAP)
    R(sl, sx2, Inches(3.35), sw2, Inches(1.5), fill=CARD)
    R(sl, sx2, Inches(3.35), sw2, Inches(0.06), fill=clr)
    T(sl, val, sx2, Inches(3.5), sw2, Inches(0.72),
      size=32, bold=True, color=clr, align=PP_ALIGN.CENTER)
    T(sl, lbl, sx2 + Inches(0.08), Inches(4.22), sw2 - Inches(0.16), Inches(0.5),
      size=11, color=T2, align=PP_ALIGN.CENTER)

T(sl, "Готов ответить на вопросы",
  MX, Inches(5.2), CW, Inches(0.42), size=16, color=T3, align=PP_ALIGN.CENTER)

pg_num(sl, 19)

# ──────────────────────────────────────────────────────────────────────────────
OUT = r"d:\Diploms\2025-2026\Musa\diploma\presentation.pptx"
prs.save(OUT)
print(f"OK  {len(prs.slides)} slides  ->  {OUT}")
